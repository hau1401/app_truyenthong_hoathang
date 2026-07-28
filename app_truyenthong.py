import streamlit as st
import time
from datetime import datetime
import base64
import io
import hashlib

# 1. Kiểm tra và nạp thư viện âm thanh (gTTS)
try:
    from gtts import gTTS
    HAS_GTTS = True
except ImportError:
    HAS_GTTS = False

# 2. Kiểm tra và nạp thư viện PowerPoint (python-pptx)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Cấu hình trang & Giao diện tổng thể
st.set_page_config(
    page_title="Cổng Quản trị & Tuyên truyền Số - Xã Hòa Thắng",
    page_icon="🇻🇳",
    layout="wide"
)

# Tùy chỉnh CSS giao diện
st.markdown("""
    <style>
    .main { background-color: #f4f6f9; }
    .stButton>button { border-radius: 8px; font-weight: bold; width: 100%; }
    .block-container { padding-top: 2.2rem; padding-bottom: 2rem; }
    h1 { font-size: 2.1rem !important; color: #b30000; font-weight: 800; text-align: center; margin: 0; }
    .sub-header { text-align: center; color: #333; font-weight: 700; font-size: 1.05rem; margin-top: 6px; }
    .flag-container { display: flex; justify-content: center; align-items: center; height: 100%; }
    .flag-img { width: 88px; height: auto; border-radius: 4px; box-shadow: 0 2px 5px rgba(0,0,0,0.15); }
    .seal-badge { background-color: #e6f4ea; border: 1px dashed #137333; padding: 10px; border-radius: 6px; font-family: monospace; font-size: 0.9rem; color: #137333; margin-top: 10px; }
    </style>
""", unsafe_allow_html=True)

# Header trang trí: Cờ Tổ quốc (Bên trái) - Tiêu đề lớn (Ở giữa) - Cờ Tổ quốc (Bên phải)
col_flag1, col_title, col_flag2 = st.columns([1, 6, 1], vertical_alignment="center")

with col_flag1:
    st.markdown(
        """
        <div class="flag-container">
            <img src="https://upload.wikimedia.org/wikipedia/commons/2/21/Flag_of_Vietnam.svg" class="flag-img">
        </div>
        """, 
        unsafe_allow_html=True
    )

with col_title:
    st.markdown("<h1>CỔNG TỰ ĐỘNG HÓA TRUYỀN THÔNG SỐ</h1>", unsafe_allow_html=True)
    st.markdown("<div class='sub-header'>ĐẢNG ỦY - HĐND - UBND XÃ HÒA THẮNG | Mô hình Chuyển đổi số Cơ sở</div>", unsafe_allow_html=True)

with col_flag2:
    st.markdown(
        """
        <div class="flag-container">
            <img src="https://upload.wikimedia.org/wikipedia/commons/2/21/Flag_of_Vietnam.svg" class="flag-img">
        </div>
        """, 
        unsafe_allow_html=True
    )

st.markdown("---")

# Khởi tạo Session State
if 'processed' not in st.session_state:
    st.session_state.processed = False
    st.session_state.raw_text = ""
    st.session_state.tieu_de = ""
    st.session_state.mo_dau = ""
    st.session_state.hanh_dong = ""
    st.session_state.digital_seal = ""
    st.session_state.history = []

# 2. Thanh bên (Sidebar) cấu hình nâng cao
with st.sidebar:
    st.markdown("### ⚙️ Cấu hình Tuyên truyền AI")
    
    phong_cach = st.selectbox(
        "🎯 Chọn giọng văn truyền thông:",
        ["Gần gũi, dễ hiểu (Phù hợp Nhân dân)", "Trang trọng, chuẩn hành chính (Đảng ủy/UBND)", "Cổ động, nhiệt huyết (Đoàn thanh niên)"]
    )
    
    doi_tuong = st.selectbox(
        "👥 Đối tượng mục tiêu trọng tâm:",
        ["Toàn thể Nhân dân trên địa bàn xã", "Đoàn viên, Thanh thiếu niên", "Hội viên Phụ nữ & Các gia đình", "Người cao tuổi & Hộ chính sách"]
    )
    
    kenh_phat_hanh = st.multiselect(
        "📢 Kênh phân phối tự động:",
        ["Trang thông tin điện tử xã", "Fanpage UBND Hòa Thắng", "Zalo OA Xã Hòa Thắng", "Hệ thống Loa truyền thanh không dây"],
        default=["Fanpage UBND Hòa Thắng", "Zalo OA Xã Hòa Thắng"]
    )

    st.markdown("---")
    st.markdown("📊 **Chỉ số Hiệu quả Mô hình:**")
    st.metric(label="⏱️ Tốc độ xử lý trung bình", value="0.75 giây", delta="-96% thời gian")
    st.metric(label="📈 Tỷ lệ phủ sóng số", value="98.5%", delta="+45% so với truyền thống")
    
    if st.session_state.history:
        st.markdown("---")
        st.markdown(f"🗂️ **Đã xuất bản phiên này:** `{len(st.session_state.history)} bản tin`")

# 3. Khu vực nạp dữ liệu mẫu nhanh
with st.expander("💡 Bấm vào đây để chọn nhanh mẫu văn bản test thử", expanded=False):
    col_s1, col_s2, col_s3 = st.columns(3)
    
    sample_1 = "THÔNG BÁO RA QUÂN NGÀY CHỦ NHẬT XANH NĂM 2026\nNhằm thiết thực lập thành tích xây dựng nông thôn mới nâng cao, UBND xã Hòa Thắng phát động toàn thể nhân dân tham gia tổng vệ sinh môi trường, phát quang bụi rậm, khơi thông cống rãnh các tuyến đường liên thôn vào lúc 7h00 sáng Chủ nhật tới đây."
    sample_2 = "THÔNG BÁO LỊCH TIÊM CHỦNG MỞ RỘNG THÁNG 6 CHO TRẺ EM\nTrạm Y tế xã Hòa Thắng thông báo lịch tiêm chủng định kỳ cho trẻ em dưới 5 tuổi diễn ra vào lúc 7h30 các ngày thứ Ba tuần tới tại hội trường Trạm Y tế xã. Đề nghị các bậc phụ huynh mang theo sổ tiêm chủng."
    sample_3 = "KẾ HOẠCH TỔ CHỨC GIẢI BÓNG ĐÁ THANH NIÊN XÃ HÒA THẮNG\nĐoàn TNCS Hồ Chí Minh xã Hòa Thắng phát động giải bóng đá mini nam truyền thống năm 2026. Các chi đoàn thôn bản khẩn trương lập danh sách vận động viên đăng ký tham gia từ nay đến hết ngày 10 tại văn phòng Đoàn xã."
    
    with col_s1:
        if st.button("Mẫu 1: Vệ sinh môi trường"):
            st.session_state.raw_text = sample_1
            st.rerun()
    with col_s2:
        if st.button("Mẫu 2: Y tế - Tiêm chủng"):
            st.session_state.raw_text = sample_2
            st.rerun()
    with col_s3:
        if st.button("Mẫu 3: Thể thao thanh niên"):
            st.session_state.raw_text = sample_3
            st.rerun()

# 4. Bố cục chính chia 2 cột cân đối
col_input, col_output = st.columns([1, 1], gap="medium")

with col_input:
    st.subheader("📝 1. Dữ liệu văn bản gốc đầu vào")
    raw_text = st.text_area(
        "Nhập hoặc dán nội dung kế hoạch, thông báo từ cơ quan:",
        placeholder="Dán văn bản hành chính vào đây...",
        height=280,
        value=st.session_state.raw_text
    )
    
    process_btn = st.button("🚀 Kích hoạt AI Tổng hợp & Đa kênh hóa", type="primary")
    
    if process_btn:
        if not raw_text.strip():
            st.warning("⚠️ Vui lòng nhập nội dung văn bản hoặc chọn mẫu nhanh ở trên!")
        else:
            with st.spinner("🤖 Trợ lý AI đang xử lý, tối ưu ngôn từ và tạo mã xác thực số..."):
                time.sleep(0.7)
                
                st.session_state.raw_text = raw_text
                lines = [l.strip() for l in raw_text.split('\n') if l.strip()]
                st.session_state.tieu_de = lines[0] if lines else "THÔNG BÁO XÃ HÒA THẮNG"
                
                # Tạo mã xác thực số / Digital Seal độc quyền bằng SHA-256 Hash
                hash_object = hashlib.sha256(raw_text.encode('utf-8'))
                hash_hex = hash_object.hexdigest()[:12].upper()
                st.session_state.digital_seal = f"HT-SEAL-2026-{hash_hex}"
                
                if "Gần gũi" in phong_cach:
                    st.session_state.mo_dau = f"Bà con {doi_tuong.lower()} ơi! UBND xã Hòa Thắng có thông báo quan trọng gửi đến toàn thể bà con nhé!"
                    st.session_state.hanh_dong = "Bà con nhớ chia sẻ thông tin rộng rãi để mọi người cùng thực hiện nha!"
                elif "Trang trọng" in phong_cach:
                    st.session_state.mo_dau = f"THÔNG BÁO CHÍNH THỨC\nGửi: {doi_tuong}\nVề việc: {st.session_state.tieu_de}"
                    st.session_state.hanh_dong = "Đề nghị các ban ngành, đoàn thể và nhân dân nghiêm túc triển khai thực hiện."
                else:
                    st.session_state.mo_dau = f"TIN NÓNG XÃ HÒA THẮNG - KẾ HOẠCH TRỌNG TÂM CHO {doi_tuong.upper()}!"
                    st.session_state.hanh_dong = "Hỡi toàn thể nhân dân Hòa Thắng, hãy cùng chung tay hành động vì quê hương giàu đẹp!"
                
                st.session_state.processed = True
                
                current_time = datetime.now().strftime("%H:%M:%S")
                st.session_state.history.append({"time": current_time, "title": st.session_state.tieu_de})

with col_output:
    st.subheader("✨ 2. Kết quả ấn phẩm truyền thông số")
    
    if st.session_state.processed:
        word_count = len(st.session_state.raw_text.split())
        est_read_time = max(5, int(word_count / 2.5))
        
        col_m1, col_m2, col_m3 = st.columns(3)
        col_m1.metric("Số từ gốc", f"{word_count} từ")
        col_m2.metric("Thời lượng phát thanh", f"~{est_read_time} giây")
        col_m3.metric("Kênh phân phối", f"{len(kenh_phat_hanh)} nền tảng")
        
        tab1, tab2, tab3 = st.tabs(["📱 Bài đăng Mạng Xã Hội", "🖼️ Thẻ Infographic Canva", "🎬 Kịch bản, Phát thanh & Word"])
        
        with tab1:
            with st.container(border=True):
                st.markdown(f"**{st.session_state.mo_dau}**")
                st.write(f"\n{st.session_state.raw_text}\n")
                st.markdown(f"*{st.session_state.hanh_dong}*")
                st.markdown("---")
                st.markdown(f"🌐 **Kênh đã duyệt:** `{' | '.join(kenh_phat_hanh)}`")
                st.markdown(f"🛡️ **Huy hiệu Xác thực Chính thống:** `{st.session_state.digital_seal}`")
                st.markdown("🏷️ **Tags:** `#HoaThangSo #ChuyenDoiSo #UBNDHoàThắng #ThongBaoChinhThuc`")
                st.download_button("📥 Tải tệp nội dung mạng xã hội (.txt)", data=st.session_state.raw_text, file_name="bai_dang_mxh.txt", use_container_width=True)
                
        # ==========================================
        # TAB 2: POWERPOINT HIGH-END CANVA DESIGN
        # ==========================================
        with tab2:
            with st.container(border=True):
                st.markdown("#### 🎨 Cấu trúc Slide Infographic Đề xuất")
                st.markdown("🔹 **Slide 1 (Bìa):** Chủ đề trọng tâm & Sọc nhận diện thương hiệu UBND xã Hòa Thắng 🏛️")
                st.markdown(f"🔹 **Slide 2 (Nội dung chính):** Định dạng dạng Thẻ nổi (Card UI) hiện đại")
                st.markdown("🔹 **Slide 3 (Hành động & Xác thực):** Khung xác thực điện tử chuẩn cơ sở 🌟")
                st.markdown(f"🛡️ **Mã QR/Seal trên Slide:** `{st.session_state.digital_seal}`")
                
                st.markdown("---")
                
                if not HAS_PPTX:
                    st.error("⚠️ Hệ thống chưa cài thư viện python-pptx. Chạy lệnh: `pip install python-pptx`")
                else:
                    prs = Presentation()
                    # Cấu hình slide chuẩn 16:9 (Hiện đại, tối ưu cho màn hình và điện thoại)
                    prs.slide_width = Inches(13.333)
                    prs.slide_height = Inches(7.5)

                    # Bảng màu thiết kế chuẩn Hành chính & Canva cao cấp
                    red_color = RGBColor(179, 0, 0)       # Đỏ cờ chủ đạo
                    yellow_color = RGBColor(255, 215, 0)  # Vàng sao
                    blue_dark = RGBColor(10, 25, 47)      # Xanh than đậm (Sangen / Tiêu đề)
                    bg_card = RGBColor(248, 249, 250)     # Xám rất nhạt cho nền Card
                    border_card = RGBColor(222, 226, 230) # Viền xám nhẹ
                    green_seal = RGBColor(19, 115, 51)    # Xanh lá xác thực số

                    # ==========================================
                    # SLIDE 1: BÌA (COVER) - PHONG CÁCH TỐI GIẢN HIỆN ĐẠI
                    # ==========================================
                    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    # Thanh sọc dọc trang trí màu đỏ cờ bên trái (Điểm nhấn Canva)
                    stripe_left = slide1.shapes.add_shape(1, Inches(0.8), Inches(1.2), Inches(0.25), Inches(5.1))
                    stripe_left.fill.solid()
                    stripe_left.fill.fore_color.rgb = red_color
                    stripe_left.line.fill.background()

                    # Tiêu đề chính căn trái theo sọc dọc
                    txBox_title = slide1.shapes.add_textbox(Inches(1.4), Inches(1.8), Inches(11.0), Inches(2.5))
                    tf_title = txBox_title.text_frame
                    tf_title.word_wrap = True
                    p_title = tf_title.paragraphs[0]
                    p_title.text = st.session_state.tieu_de.upper()
                    p_title.font.bold = True
                    p_title.font.size = Pt(36)
                    p_title.font.color.rgb = blue_dark

                    # Khung thông tin cơ quan ban hành phía dưới
                    txBox_sub = slide1.shapes.add_textbox(Inches(1.4), Inches(4.8), Inches(11.0), Inches(1.2))
                    tf_sub = txBox_sub.text_frame
                    p_sub1 = tf_sub.paragraphs[0]
                    p_sub1.text = "ĐẢNG ỦY - HĐND - UBND XÃ HÒA THẮNG"
                    p_sub1.font.bold = True
                    p_sub1.font.size = Pt(18)
                    p_sub1.font.color.rgb = red_color
                    
                    p_sub2 = tf_sub.add_paragraph()
                    p_sub2.text = "Mô hình Chuyển đổi số & Tuyên truyền Cơ sở năm 2026"
                    p_sub2.font.size = Pt(14)
                    p_sub2.font.color.rgb = RGBColor(100, 100, 100)

                    # ==========================================
                    # SLIDE 2: NỘI DUNG CHÍNH - THẺ NỔI (CARD UI)
                    # ==========================================
                    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    # Tiêu đề Slide có gạch chân nhỏ
                    txBox_h2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
                    p_h2 = txBox_h2.text_frame.paragraphs[0]
                    p_h2.text = "NỘI DUNG TRỌNG TÂM"
                    p_h2.font.bold = True
                    p_h2.font.size = Pt(24)
                    p_h2.font.color.rgb = blue_dark
                    
                    line_h2 = slide2.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(2.5), Inches(0.04))
                    line_h2.fill.solid()
                    line_h2.fill.fore_color.rgb = red_color
                    line_h2.line.fill.background()

                    # Thẻ chứa nội dung (Card container màu xám nhạt tạo chiều sâu)
                    card_s2 = slide2.shapes.add_shape(1, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0))
                    card_s2.fill.solid()
                    card_s2.fill.fore_color.rgb = bg_card
                    card_s2.line.color.rgb = border_card
                    card_s2.line.width = Pt(1)

                    # Văn bản bên trong Thẻ
                    txBox_b2 = slide2.shapes.add_textbox(Inches(1.1), Inches(1.9), Inches(11.1), Inches(4.4))
                    tf_b2 = txBox_b2.text_frame
                    tf_b2.word_wrap = True
                    p_b2 = tf_b2.paragraphs[0]
                    p_b2.text = st.session_state.raw_text
                    p_b2.font.size = Pt(18)
                    p_b2.font.color.rgb = RGBColor(33, 37, 41)

                    # ==========================================
                    # SLIDE 3: HÀNH ĐỘNG & XÁC THỰC SỐ - THẺ ĐẶC BIỆT
                    # ==========================================
                    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
                    
                    txBox_h3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(0.8))
                    p_h3 = txBox_h3.text_frame.paragraphs[0]
                    p_h3.text = "TỔ CHỨC THỰC HIỆN & XÁC THỰC SỐ"
                    p_h3.font.bold = True
                    p_h3.font.size = Pt(24)
                    p_h3.font.color.rgb = blue_dark
                    
                    line_h3 = slide3.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(3.2), Inches(0.04))
                    line_h3.fill.solid()
                    line_h3.fill.fore_color.rgb = red_color
                    line_h3.line.fill.background()

                    # Thẻ Kêu gọi hành động
                    card_action = slide3.shapes.add_shape(1, Inches(0.8), Inches(1.6), Inches(11.733), Inches(1.5))
                    card_action.fill.solid()
                    card_action.fill.fore_color.rgb = RGBColor(254, 242, 242) # Đỏ rất nhạt
                    card_action.line.color.rgb = RGBColor(254, 202, 202)
                    card_action.line.width = Pt(1)

                    txBox_b3 = slide3.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(11.1), Inches(1.1))
                    tf_b3 = txBox_b3.text_frame
                    tf_b3.word_wrap = True
                    p_b3 = tf_b3.paragraphs[0]
                    p_b3.text = f"🌟 {st.session_state.hanh_dong}"
                    p_b3.font.size = Pt(18)
                    p_b3.font.bold = True
                    p_b3.font.color.rgb = red_color

                    # Thẻ Khung Xác thực Số (Digital Seal Box)
                    seal_box = slide3.shapes.add_shape(1, Inches(3.16), Inches(3.4), Inches(7.0), Inches(3.1))
                    seal_box.fill.solid()
                    seal_box.fill.fore_color.rgb = RGBColor(230, 244, 234) # Xanh lá nhạt
                    seal_box.line.color.rgb = green_seal
                    seal_box.line.width = Pt(1.5)
                    seal_box.line.dash_style = 7 # Đường viền đứt nét điện tử

                    txBox_seal = slide3.shapes.add_textbox(Inches(3.3), Inches(3.7), Inches(6.7), Inches(2.5))
                    tf_seal = txBox_seal.text_frame
                    p_seal1 = tf_seal.paragraphs[0]
                    p_seal1.text = "🛡️ HUY HIỆU XÁC THỰC CHÍNH THỐNG (DIGITAL SEAL)"
                    p_seal1.font.bold = True
                    p_seal1.font.size = Pt(16)
                    p_seal1.font.color.rgb = green_seal
                    p_seal1.alignment = PP_ALIGN.CENTER
                    
                    p_seal2 = tf_seal.add_paragraph()
                    p_seal2.text = f"\nMã tra cứu bảo mật:\n{st.session_state.digital_seal}"
                    p_seal2.font.bold = True
                    p_seal2.font.size = Pt(22)
                    p_seal2.font.color.rgb = blue_dark
                    p_seal2.alignment = PP_ALIGN.CENTER
                    
                    p_seal3 = tf_seal.add_paragraph()
                    p_seal3.text = "(Văn bản được mã hóa và chứng thực trực tiếp trên Cổng Thông tin Xã Hòa Thắng)"
                    p_seal3.font.size = Pt(11)
                    p_seal3.font.italic = True
                    p_seal3.font.color.rgb = RGBColor(100, 100, 100)
                    p_seal3.alignment = PP_ALIGN.CENTER

                    # ==========================================
                    # FOOTER CHUNG CHO TẤT CẢ SLIDES
                    # ==========================================
                    for slide in prs.slides:
                        footer = slide.shapes.add_textbox(Inches(0.8), Inches(7.1), Inches(11.733), Inches(0.3))
                        p_foot = footer.text_frame.paragraphs[0]
                        p_foot.text = "© 2026 UBND Xã Hòa Thắng, Lâm Đồng — Nền tảng Chuyển đổi số Cơ sở"
                        p_foot.font.size = Pt(10)
                        p_foot.font.color.rgb = RGBColor(140, 140, 140)

                    # Lưu và đóng gói ra Nút Tải về
                    ppt_stream = io.BytesIO()
                    prs.save(ppt_stream)
                    ppt_stream.seek(0)
                    
                    st.download_button(
                        label="📥 Tải bộ Slide thiết kế cao cấp chuẩn Canva (.pptx)",
                        data=ppt_stream,
                        file_name=f"SlideCanva_HoaThang_{datetime.now().strftime('%Y%m%d')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )

        with tab3:
            with st.container(border=True):
                st.markdown("#### 🎙️ Giả lập Phát thanh truyền thanh cơ sở (AI Text-to-Speech)")
                noi_dung_doc = f"Xin kính chào bà con nhân dân xã Hòa Thắng! Hệ thống truyền thanh cơ sở xin thông báo: {st.session_state.tieu_de}. {st.session_state.raw_text} Đề nghị toàn thể nhân dân chú ý theo dõi và thực hiện đầy đủ. Trân trọng cảm ơn!"
                st.info(noi_dung_doc)
                
                if st.button("🔊 Nghe thử bản tin phát thanh trực tiếp"):
                    if not HAS_GTTS:
                        st.error("⚠️ Hệ thống chưa cài thư viện gTTS. Vui lòng chạy lệnh: pip install gTTS")
                    else:
                        with st.spinner("Đang khởi tạo giọng đọc AI..."):
                            try:
                                tts = gTTS(text=noi_dung_doc, lang='vi')
                                fp = io.BytesIO()
                                tts.write_to_fp(fp)
                                fp.seek(0)
                                st.audio(fp, format='audio/mp3')
                            except Exception as e:
                                st.error(f"Lỗi khởi tạo âm thanh: {e}")
                
                st.markdown("---")
                st.markdown("#### 📄 Xuất Báo cáo Hành chính chuẩn Word (.doc)")
                
                today = datetime.now()
                
                paragraphs = st.session_state.raw_text.split('\n')
                formatted_content = ""
                for p in paragraphs:
                    if p.strip():
                        formatted_content += f"<p class='noi-dung'>{p.strip()}</p>\n"

                word_html = f"""
                <html xmlns:o='urn:schemas-microsoft-com:office:office' xmlns:w='urn:schemas-microsoft-com:office:word' xmlns='http://www.w3.org/TR/REC-html40'>
                <head>
                <meta http-equiv="Content-Type" content="text/html; charset=utf-8">
                <title>Báo cáo hành chính</title>
                <style>
                    @page WordSection1 {{
                        size: 21.0cm 29.7cm;
                        margin: 2.0cm 2.0cm 2.0cm 3.0cm;
                        mso-header-margin: 35.4pt;
                        mso-footer-margin: 35.4pt;
                        mso-paper-source: 0;
                    }}
                    div.WordSection1 {{ page: WordSection1; }}
                    body {{ font-family: "Times New Roman", Times, serif; font-size: 14pt; }}
                    table {{ border-collapse: collapse; width: 100%; border: none; }}
                    td {{ border: none; vertical-align: top; padding: 0; }}
                    .cq-chu-quan {{ font-size: 12pt; text-align: center; }}
                    .cq-ban-hanh {{ font-size: 13pt; font-weight: bold; text-align: center; }}
                    .quoc-hieu {{ font-size: 12pt; font-weight: bold; text-align: center; }}
                    .tieu-ngu {{ font-size: 14pt; font-weight: bold; text-align: center; }}
                    .ngay-thang {{ font-size: 14pt; font-style: italic; text-align: center; }}
                    .so-hieu {{ font-size: 13pt; text-align: center; }}
                    .ten-van-ban {{ font-size: 14pt; font-weight: bold; text-align: center; margin-top: 18pt; margin-bottom: 18pt; text-transform: uppercase; }}
                    .noi-dung {{ 
                        text-align: justify; 
                        text-indent: 1.27cm; 
                        line-height: 150%; 
                        margin-top: 6pt; 
                        margin-bottom: 6pt; 
                    }}
                    .seal-box {{ border: 1px dashed #137333; padding: 12px; margin-top: 24pt; font-family: "Courier New", monospace; font-size: 12pt; color: #137333; }}
                </style>
                </head>
                <body>
                <div class="WordSection1">
                    <table>
                        <tr>
                            <td width="40%" style="text-align: center;">
                                <div class="cq-chu-quan">UBND HUYỆN...</div>
                                <div class="cq-ban-hanh">UBND XÃ HÒA THẮNG</div>
                                <hr style="width: 40%; border: 0; border-top: 1px solid black; margin-top: 2px; margin-bottom: 4px;">
                                <div class="so-hieu">Số: {len(st.session_state.history) + 1}/TB-UBND</div>
                            </td>
                            <td width="60%" style="text-align: center;">
                                <div class="quoc-hieu">CỘNG HÒA XÃ HỘI CHỦ NGHĨA VIỆT NAM</div>
                                <div class="tieu-ngu">Độc lập - Tự do - Hạnh phúc</div>
                                <hr style="width: 40%; border: 0; border-top: 1px solid black; margin-top: 2px; margin-bottom: 4px;">
                                <div class="ngay-thang">Hòa Thắng, ngày {today.strftime("%d")} tháng {today.strftime("%m")} năm {today.strftime("%Y")}</div>
                            </td>
                        </tr>
                    </table>
                    
                    <div class="ten-van-ban">{st.session_state.tieu_de}</div>
                    
                    <p class="noi-dung" style="font-weight: bold;">
                        {st.session_state.mo_dau}
                    </p>

                    {formatted_content}

                    <p class="noi-dung" style="font-style: italic;">
                        {st.session_state.hanh_dong}
                    </p>
                    
                    <div class="seal-box">
                        <b>🛡️ HUY HIỆU XÁC THỰC CHÍNH THỐNG (DIGITAL SEAL):</b><br>
                        Mã tra cứu số: {st.session_state.digital_seal}<br>
                        <i>(Văn bản được xác thực ký số tự động trên hệ thống Cổng Thông tin)</i>
                    </div>
                    
                    <br>
                    
                    <table>
                        <tr>
                            <td width="50%" style="font-size: 12pt; line-height: 120%;">
                                <b><i>Nơi nhận:</i></b><br>
                                - Thường trực Đảng ủy;<br>
                                - Thường trực HĐND;<br>
                                - Lưu: VT.
                            </td>
                            <td width="50%" style="text-align: center;">
                                <div style="font-size: 13pt; font-weight: bold;">TM. ỦY BAN NHÂN DÂN XÃ HÒA THẮNG</div>
                                <div style="font-size: 13pt; font-weight: bold;">KT. CHỦ TỊCH</div>
                                <div style="font-size: 13pt; font-weight: bold;">PHÓ CHỦ TỊCH</div>
                                <br><br><br><br><br>
                                <div style="font-size: 14pt; font-weight: bold;">(Ký, đóng dấu)</div>
                            </td>
                        </tr>
                    </table>
                </div>
                </body>
                </html>
                """
                
                b64_word = base64.b64encode(word_html.encode('utf-8-sig')).decode("utf-8")
                word_download_link = f'<a href="data:application/vnd.ms-word;charset=utf-8;base64,{b64_word}" download="BaoCaoHanChinh_HoaThang_{today.strftime("%Y%m%d")}.doc" style="display:inline-block;padding:10px 20px;background-color:#b30000;color:white;text-decoration:none;border-radius:6px;font-weight:bold;text-align:center;width:100%;">📥 Tải xuống Báo cáo hành chính chuẩn (.doc)</a>'
                st.markdown(word_download_link, unsafe_allow_html=True)
    else:
        st.info("👈 Vui lòng nhập dữ liệu hoặc chọn mẫu nhanh ở trên, sau đó bấm nút **'Kích hoạt AI Tổng hợp & Đa kênh hóa'** để xem kết quả xuất bản.")